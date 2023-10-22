from pptx import Presentation
from pptx.dml.color import RGBColor
import json

def create_presentation(input_json_path, template_name):
    output_pptx_path = "NewModified_Presentation.pptx"
    # Construct the path to the template presentation
    presentation_path = f'Static/{template_name}.pptx'
    print(presentation_path)
    # Load the existing presentation template
    presentation = Presentation(presentation_path)

    slides_data = input_json_path
    print(slides_data)
    # Loop through each slide in the presentation and the JSON data
    for i, (slide, slide_data) in enumerate(zip(presentation.slides, slides_data["slides"])):
        # First Slide: Title Slide
        if i == 0:
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]
            if slide.shapes.placeholders[1]:
                slide.shapes.placeholders[1].text = slide_data["subtitle"]
        # Second Slide: Table of Contents
        elif i == 1:
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]
            if slide.shapes.placeholders[1]:
                slide.shapes.placeholders[1].text = "\n".join(slide_data["bullets"])
        # Last Slide: Thank You
        elif i == len(slides_data["slides"]) - 1:
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]
        # In-between Slides: Content Slides
        else:
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]
            if slide.shapes.placeholders[1]:
                slide.shapes.placeholders[1].text = "\n".join(slide_data["content"])

    # Save the modified presentation
    presentation.save(output_pptx_path)


jsonVal = {
    "slides": [
      {
        "title": "Dolphins: The Enigmatic Sea Creatures",
        "subtitle": "Exploring the Wonders of the Ocean"
      },
      {
        "title": "Table of Contents",
        "bullets": [
          "Introduction to Dolphins",
          "Dolphin Species",
          "Dolphin Behavior",
          "Dolphin Intelligence",
          "Dolphin Habitats",
          "Dolphin Conservation",
          "Fun Facts about Dolphins",
          "Thank You"
        ]
      },
      {
        "title": "Introduction to Dolphins",
        "content": [
          "Dolphins, often known as the ocean's ambassadors, are magnificent marine mammals.",
          "In this presentation, we'll embark on a journey to discover their extraordinary world.",
          "Join us as we dive into the enchanting realm of dolphins."
        ]
      },
      {
        "title": "Dolphin Species",
        "content": [
          "The world of dolphins is diverse, with various species inhabiting our oceans and rivers.",
          "We will explore the most common species, from the playful Bottlenose Dolphin to the majestic Orca.",
          "Get ready to meet the amazing dolphin species of the world."
        ]
      },
      {
        "title": "Dolphin Behavior",
        "content": [
          "Dolphins are not just graceful swimmers; they exhibit complex social behaviors and communication.",
          "Join us as we delve into their intriguing world of clicks, whistles, and group dynamics.",
          "Discover the secrets behind their fascinating behaviors."
        ]
      },
      {
        "title": "Dolphin Intelligence",
        "content": [
          "Dolphins are renowned for their intelligence and problem-solving abilities.",
          "We'll explore their remarkable cognitive skills, including self-awareness and innovative thinking.",
          "Prepare to be amazed by the genius of these marine mammals."
        ]
      },
      {
        "title": "Dolphin Habitats",
        "content": [
          "Dolphins inhabit a wide range of environments, from coastal waters to the deep ocean.",
          "We'll take a journey to their various habitats and the unique adaptations that help them thrive.",
          "Discover the significance of their homes in the underwater world."
        ]
      },
      {
        "title": "Dolphin Conservation",
        "content": [
          "Conservation efforts are vital to protect dolphins from the many threats they face, such as pollution and habitat loss.",
          "Learn about the initiatives, research, and marine protected areas dedicated to their well-being.",
          "Join us in the mission to preserve these remarkable creatures."
        ]
      },
      {
        "title": "Fun Facts about Dolphins",
        "content": [
          "Dolphins are full of surprises! Did you know they can swim at speeds of up to 60 kilometers per hour?",
          "We'll also share fascinating tidbits about their communication, echolocation, and playful behaviors.",
          "Get ready for a splash of fun facts about dolphins."
        ]
      },
      {
        "title": "Thank You"
      }
    ]
}
# Usage example
#create_presentation(jsonVal, 1)

